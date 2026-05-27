from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from thucthengay.config import load_project_config


def write_json(path: Path, data: dict[str, object]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(data), encoding="utf-8")


def write_template_pptx(
    path: Path,
    *,
    slide_count: int = 1,
    width: int | None = None,
    height: int | None = None,
    shape_count: int = 1,
) -> int:
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]
    shape_width = Inches(4) if width is None else width
    shape_height = Inches(3) if height is None else height
    first_shape_id = 0
    for index in range(slide_count):
        slide = presentation.slides.add_slide(blank_layout)
        for shape_index in range(shape_count):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1 + shape_index),
                Inches(1 + shape_index),
                shape_width,
                shape_height,
            )
            if index == 0 and shape_index == 0:
                first_shape_id = int(shape.shape_id)
    presentation.save(path)
    return first_shape_id


def target_config(
    target_id: str,
    sort_order: int,
    *,
    enabled: bool = True,
    template_pptx_file: str | None = None,
    map_element_id: int = 2,
) -> dict[str, object]:
    return {
        "id": target_id,
        "enabled": enabled,
        "sort_order": sort_order,
        "name": target_id,
        "geojson_file": f"targets/{target_id}.geojson",
        "coordinate": [106.7, 10.8],
        "scale": 50000,
        "grid": {"interval": {"minutes": 1}},
        "export": {
            "template_pptx_file": template_pptx_file or f"templates/{target_id}.pptx",
            "placeholders": [
                {"field": "map_image", "kind": "map_image", "element_id": map_element_id}
            ],
        },
    }


def prepare_target_files(root: Path, target_id: str, template_name: str | None = None) -> int:
    (root / "targets").mkdir(parents=True, exist_ok=True)
    (root / "targets" / f"{target_id}.geojson").write_text("{}", encoding="utf-8")
    return write_template_pptx(root / "templates" / (template_name or f"{target_id}.pptx"))


def test_load_project_config_filters_enabled_targets_and_resolves_paths(tmp_path: Path) -> None:
    id_b = prepare_target_files(tmp_path, "target_b")
    id_a = prepare_target_files(tmp_path, "target_a")
    write_json(
        tmp_path / "config.json",
        {
            "targets": [
                target_config("target_b", 2, map_element_id=id_b),
                target_config("disabled", 0, enabled=False),
                target_config("target_a", 1, map_element_id=id_a),
            ]
        },
    )

    result = load_project_config(tmp_path / "config.json")

    assert result.ok is True
    assert [target.id for target in result.enabled_targets] == ["target_a", "target_b"]
    assert result.target_paths["target_a"].geojson_file == (
        tmp_path / "targets" / "target_a.geojson"
    ).resolve()
    assert result.target_paths["target_a"].template_pptx_file == (
        tmp_path / "templates" / "target_a.pptx"
    ).resolve()
    assert result.template_metadata["target_a"].map_frame.width == 288
    assert result.enabled_targets[0].metadata["template_metadata"]["template_pptx"] == str(
        (tmp_path / "templates" / "target_a.pptx").resolve()
    )


def test_disabled_targets_are_not_schema_or_reference_blockers(tmp_path: Path) -> None:
    element_id = prepare_target_files(tmp_path, "target_a")
    disabled_invalid_target = {
        "id": "disabled",
        "enabled": False,
        "sort_order": 0,
        "name": "Disabled",
    }
    write_json(
        tmp_path / "config.json",
        {
            "targets": [
                disabled_invalid_target,
                target_config("target_a", 1, map_element_id=element_id),
            ]
        },
    )

    result = load_project_config(tmp_path / "config.json")

    assert result.ok is True
    assert [target.id for target in result.enabled_targets] == ["target_a"]
    assert "disabled" not in result.target_paths


def test_invalid_config_returns_vietnamese_field_path_issue(tmp_path: Path) -> None:
    data = target_config("target_a", 1)
    data["coordinate"] = [999, 10]
    write_json(tmp_path / "config.json", {"targets": [data]})

    result = load_project_config(tmp_path / "config.json")

    assert result.ok is False
    assert result.issues[0].blocking is True
    assert "`targets.0.coordinate`" in result.issues[0].message
    assert "[lon, lat]" in (result.issues[0].remediation or "")


def test_missing_references_create_blocking_target_issues(tmp_path: Path) -> None:
    write_json(tmp_path / "config.json", {"targets": [target_config("target_a", 1)]})

    result = load_project_config(tmp_path / "config.json")
    issue_ids = {issue.issue_id for issue in result.issues}

    assert result.ok is False
    assert "target.geojson_missing" in issue_ids
    assert "target.template_pptx_missing" in issue_ids
    assert all(issue.blocking for issue in result.issues)


def test_config_directory_path_returns_structured_issue(tmp_path: Path) -> None:
    result = load_project_config(tmp_path)

    assert result.ok is False
    assert result.issues[0].issue_id == "config.file_unreadable"
    assert result.issues[0].blocking is True


def test_template_pptx_resolves_relative_to_config_file(tmp_path: Path) -> None:
    (tmp_path / "targets").mkdir(parents=True)
    (tmp_path / "targets" / "target_a.geojson").write_text("{}", encoding="utf-8")
    element_id = write_template_pptx(tmp_path / "templates" / "nested" / "target_a.pptx")
    write_json(
        tmp_path / "config.json",
        {
            "targets": [
                target_config(
                    "target_a",
                    1,
                    template_pptx_file="templates/nested/target_a.pptx",
                    map_element_id=element_id,
                )
            ]
        },
    )

    result = load_project_config(tmp_path / "config.json")

    assert result.ok is True
    assert result.target_paths["target_a"].template_pptx_file == (
        tmp_path / "templates" / "nested" / "target_a.pptx"
    ).resolve()
    assert result.template_metadata["target_a"].template_pptx == str(
        (tmp_path / "templates" / "nested" / "target_a.pptx").resolve()
    )


def test_template_pptx_must_have_exactly_one_slide(tmp_path: Path) -> None:
    (tmp_path / "targets").mkdir(parents=True)
    (tmp_path / "targets" / "target_a.geojson").write_text("{}", encoding="utf-8")
    element_id = write_template_pptx(tmp_path / "templates" / "target_a.pptx", slide_count=2)
    write_json(
        tmp_path / "config.json",
        {"targets": [target_config("target_a", 1, map_element_id=element_id)]},
    )

    result = load_project_config(tmp_path / "config.json")

    assert result.ok is False
    assert {issue.issue_id for issue in result.issues} == {
        "target.template_pptx_slide_count_invalid"
    }


def test_required_element_id_must_exist_in_template(tmp_path: Path) -> None:
    prepare_target_files(tmp_path, "target_a")
    write_json(
        tmp_path / "config.json",
        {"targets": [target_config("target_a", 1, map_element_id=999)]},
    )

    result = load_project_config(tmp_path / "config.json")

    assert result.ok is False
    assert {issue.issue_id for issue in result.issues} == {"target.template_element_missing"}


def test_invalid_map_frame_returns_structured_template_issue(tmp_path: Path) -> None:
    (tmp_path / "targets").mkdir(parents=True)
    (tmp_path / "targets" / "target_a.geojson").write_text("{}", encoding="utf-8")
    element_id = write_template_pptx(tmp_path / "templates" / "target_a.pptx", width=0)
    write_json(
        tmp_path / "config.json",
        {"targets": [target_config("target_a", 1, map_element_id=element_id)]},
    )

    result = load_project_config(tmp_path / "config.json")

    assert result.ok is False
    assert {issue.issue_id for issue in result.issues} == {"target.template_map_frame_invalid"}


def test_duplicate_placeholder_element_ids_are_blocking(tmp_path: Path) -> None:
    element_id = prepare_target_files(tmp_path, "target_a")
    target = target_config("target_a", 1, map_element_id=element_id)
    target["export"]["placeholders"].append(  # type: ignore[index, union-attr]
        {"field": "target_title", "kind": "text", "element_id": element_id}
    )
    write_json(tmp_path / "config.json", {"targets": [target]})

    result = load_project_config(tmp_path / "config.json")

    assert result.ok is False
    assert {issue.issue_id for issue in result.issues} == {"target.template_element_duplicate"}


def test_multiple_required_map_placeholders_are_blocking(tmp_path: Path) -> None:
    (tmp_path / "targets").mkdir(parents=True)
    (tmp_path / "targets" / "target_a.geojson").write_text("{}", encoding="utf-8")
    first_id = write_template_pptx(
        tmp_path / "templates" / "target_a.pptx",
        shape_count=2,
    )
    second_id = first_id + 1
    target = target_config("target_a", 1, map_element_id=first_id)
    target["export"]["placeholders"].append(  # type: ignore[index, union-attr]
        {"field": "map_image_secondary", "kind": "map_image", "element_id": second_id}
    )
    write_json(tmp_path / "config.json", {"targets": [target]})

    result = load_project_config(tmp_path / "config.json")

    assert result.ok is False
    assert {issue.issue_id for issue in result.issues} == {"target.template_map_element_ambiguous"}


def test_multiple_template_files_surface_compatibility_warning(tmp_path: Path) -> None:
    id_a = prepare_target_files(tmp_path, "target_a")
    id_b = prepare_target_files(tmp_path, "target_b")
    write_json(
        tmp_path / "config.json",
        {
            "targets": [
                target_config("target_a", 1, map_element_id=id_a),
                target_config("target_b", 2, map_element_id=id_b),
            ]
        },
    )

    result = load_project_config(tmp_path / "config.json")

    assert result.ok is True
    issue = next(
        issue
        for issue in result.issues
        if issue.issue_id == "target.template_compatibility_unknown"
    )
    assert "target_a" in (issue.remediation or "")
    assert "target_b" in (issue.remediation or "")
