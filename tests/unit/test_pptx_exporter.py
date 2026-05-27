from __future__ import annotations

from datetime import date, time
from pathlib import Path

import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from thucthengay.export import ensure_final_renders_for_export, export_combined_pptx
from thucthengay.models import (
    Composition,
    GridConfig,
    GridInterval,
    ImageLayer,
    MapFrame,
    MetadataStatus,
    PlaceholderType,
    TargetConfig,
    TemplateMetadata,
    TemplatePlaceholder,
    ViewState,
)
from thucthengay.render import RasterRenderResult, RenderSpec
from thucthengay.workspace import WorkspaceService


def _write_template(path: Path, *, title: str = "Template") -> tuple[int, int]:
    path.parent.mkdir(parents=True, exist_ok=True)
    logo_path = path.with_suffix(".logo.png")
    Image.new("RGB", (16, 16), (200, 40, 40)).save(logo_path)
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5.625)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    map_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5),
        Inches(0.75),
        Inches(4),
        Inches(3),
    )
    text_shape = slide.shapes.add_textbox(
        Inches(5),
        Inches(0.75),
        Inches(4),
        Inches(0.5),
    )
    text_shape.text = title
    slide.shapes.add_picture(str(logo_path), Inches(9), Inches(0.25), width=Inches(0.4))
    presentation.save(path)
    return int(map_shape.shape_id), int(text_shape.shape_id)


def _target(
    template_path: Path,
    map_element_id: int,
    text_element_id: int,
    *,
    target_id: str = "alpha",
    text_field: str = "target_title",
    text_required: bool = True,
) -> TargetConfig:
    placeholders = [
        TemplatePlaceholder(
            field="map",
            element_id=map_element_id,
            kind=PlaceholderType.MAP_IMAGE,
            required=True,
        ),
        TemplatePlaceholder(
            field=text_field,
            element_id=text_element_id,
            kind=PlaceholderType.TEXT,
            required=text_required,
        ),
    ]
    return TargetConfig(
        id=target_id,
        name=f"{target_id.title()} Name",
        alias=target_id.upper(),
        title=f"{target_id.title()} Title",
        geojson_file=f"targets/{target_id}.geojson",
        coordinate=[106.7, 10.8],
        scale=50000,
        grid=GridConfig(interval=GridInterval(minutes=1)),
        export={
            "template_pptx_file": str(template_path),
            "txt_line_template": "{slide_number}|{target_id}|{capture_date}|{time_label}",
            "placeholders": [
                item.model_dump(mode="json")
                for item in placeholders
            ],
        },
        metadata={
            "template_metadata": TemplateMetadata(
                template_pptx=str(template_path),
                slide_index=0,
                map_frame=MapFrame(x=36, y=54, width=288, height=216),
                placeholders=placeholders,
            ).model_dump(mode="json")
        },
    )


def _composition(
    composition_id: str,
    *,
    target_id: str = "alpha",
    capture_date: date = date(2026, 5, 25),
    review_order: int = 1,
) -> Composition:
    return Composition(
        composition_id=composition_id,
        target_id=target_id,
        capture_date=capture_date,
        view=ViewState(center=[106.7, 10.8], scale=50000),
        reviewed=True,
        ready=True,
        include=True,
        needs_revalidation=False,
        review_order=review_order,
        layers=[
            ImageLayer(
                layer_id=f"{composition_id}-layer",
                source_path=f"{composition_id}.tif",
                cache_path=f"cache/{composition_id}.tif",
                order=0,
                visible=True,
                capture_date=capture_date,
                capture_time=time(8, 30),
                metadata_status=MetadataStatus.VALID,
            )
        ],
    )


def _workspace(tmp_path: Path, *compositions: Composition) -> WorkspaceService:
    service = WorkspaceService(tmp_path / "workspace")
    service.initialize(config_path="config.json")
    for composition in compositions:
        service.write_composition(composition)
    return service


def _success_render(spec: RenderSpec, is_cancelled=None) -> RasterRenderResult:
    return RasterRenderResult(
        canvas=np.full((spec.output_height, spec.output_width, 3), 90, dtype=np.uint8),
        painted_layer_ids=tuple(layer.layer_id for layer in spec.visible_layers),
    )


def test_export_combined_pptx_vertical_slice_replaces_map_and_text(tmp_path: Path) -> None:
    map_id, text_id = _write_template(tmp_path / "templates" / "alpha.pptx")
    target = _target(tmp_path / "templates" / "alpha.pptx", map_id, text_id)
    service = _workspace(tmp_path, _composition("alpha__20260525"))
    ensure_final_renders_for_export(service, [target], render=_success_render)

    output_path = service.paths.exports / "report.pptx"
    result = export_combined_pptx(service, [target], output_path=output_path)

    assert result.ok is True
    assert result.summary.slide_count == 1
    assert result.pptx_path == "exports/report.pptx"
    assert result.exported[0].composition_id == "alpha__20260525"
    presentation = Presentation(str(output_path))
    assert len(presentation.slides) == 1
    slide = presentation.slides[0]
    assert any(shape.shape_type == 13 for shape in slide.shapes)
    slide_text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
    assert "Alpha Title" in slide_text


def test_export_combined_pptx_orders_slides_by_review_order(tmp_path: Path) -> None:
    map_id, text_id = _write_template(tmp_path / "templates" / "alpha.pptx")
    target = _target(tmp_path / "templates" / "alpha.pptx", map_id, text_id)
    service = _workspace(
        tmp_path,
        _composition("alpha__20260526", capture_date=date(2026, 5, 26), review_order=2),
        _composition("alpha__20260525", capture_date=date(2026, 5, 25), review_order=1),
    )
    ensure_final_renders_for_export(service, [target], render=_success_render)

    result = export_combined_pptx(
        service,
        [target],
        output_path=service.paths.exports / "ordered.pptx",
    )

    assert result.ok is True
    assert [row.composition_id for row in result.exported] == [
        "alpha__20260525",
        "alpha__20260526",
    ]
    assert [row.slide_number for row in result.exported] == [1, 2]


def test_export_combined_pptx_blocks_unresolved_required_text_placeholder(
    tmp_path: Path,
) -> None:
    map_id, text_id = _write_template(tmp_path / "templates" / "alpha.pptx")
    target = _target(
        tmp_path / "templates" / "alpha.pptx",
        map_id,
        text_id,
        text_field="unsupported_required_value",
    )
    service = _workspace(tmp_path, _composition("alpha__20260525"))
    ensure_final_renders_for_export(service, [target], render=_success_render)

    result = export_combined_pptx(
        service,
        [target],
        output_path=service.paths.exports / "blocked.pptx",
    )

    assert result.ok is False
    assert "export.pptx_placeholder_unresolved" in {issue.issue_id for issue in result.issues}
    assert not (service.paths.exports / "blocked.pptx").exists()


def test_export_combined_pptx_blocks_missing_final_png(tmp_path: Path) -> None:
    map_id, text_id = _write_template(tmp_path / "templates" / "alpha.pptx")
    target = _target(tmp_path / "templates" / "alpha.pptx", map_id, text_id)
    service = _workspace(tmp_path, _composition("alpha__20260525"))

    result = export_combined_pptx(
        service,
        [target],
        output_path=service.paths.exports / "missing-render.pptx",
    )

    assert result.ok is False
    assert "export.final_render_missing" in {issue.issue_id for issue in result.issues}
    assert not (service.paths.exports / "missing-render.pptx").exists()


def test_export_combined_pptx_rejects_out_of_workspace_output_path(tmp_path: Path) -> None:
    map_id, text_id = _write_template(tmp_path / "templates" / "alpha.pptx")
    target = _target(tmp_path / "templates" / "alpha.pptx", map_id, text_id)
    service = _workspace(tmp_path, _composition("alpha__20260525"))
    ensure_final_renders_for_export(service, [target], render=_success_render)
    output_path = tmp_path / "outside.pptx"

    result = export_combined_pptx(service, [target], output_path=output_path)

    assert result.ok is False
    assert "export.pptx_write_failed" in {issue.issue_id for issue in result.issues}
    assert not output_path.exists()
