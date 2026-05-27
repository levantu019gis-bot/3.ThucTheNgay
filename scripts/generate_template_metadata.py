#!/usr/bin/env python3
"""Generate architecture-compliant PPTX template metadata."""

from __future__ import annotations

import argparse
import io
import json
import os
import sys
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET
from zipfile import ZipFile

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PROJECT_ROOT = Path(__file__).resolve().parents[1]
SRC_ROOT = PROJECT_ROOT / "src"
if str(SRC_ROOT) not in sys.path:
    sys.path.insert(0, str(SRC_ROOT))

from thucthengay.models import (  # noqa: E402
    MapFrame,
    TemplateMetadata,
    TemplatePlaceholder,
)

EMU_PER_INCH = 914400
NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
}


def emu_to_inches(value: int) -> float:
    return round(value / EMU_PER_INCH, 4)


def enum_name(value: Any) -> str | None:
    if value is None:
        return None
    return getattr(value, "name", str(value))


def color_payload(color: Any) -> dict[str, Any]:
    result: dict[str, Any] = {}
    try:
        if color.rgb is not None:
            result["rgb"] = str(color.rgb)
    except Exception:
        pass
    try:
        if color.theme_color is not None:
            result["theme_color"] = enum_name(color.theme_color)
    except Exception:
        pass
    try:
        result["brightness"] = color.brightness
    except Exception:
        pass
    return result


def fill_payload(shape: Any) -> dict[str, Any] | None:
    fill = getattr(shape, "fill", None)
    if fill is None:
        return None
    result: dict[str, Any] = {"type": enum_name(getattr(fill, "type", None))}
    try:
        result["fore_color"] = color_payload(fill.fore_color)
    except Exception:
        pass
    try:
        result["back_color"] = color_payload(fill.back_color)
    except Exception:
        pass
    return result


def line_payload(shape: Any) -> dict[str, Any] | None:
    line = getattr(shape, "line", None)
    if line is None:
        return None
    result: dict[str, Any] = {}
    for attr in ("width", "dash_style", "begin_arrowhead", "end_arrowhead"):
        try:
            result[attr] = enum_name(getattr(line, attr))
        except Exception:
            pass
    try:
        result["color"] = color_payload(line.color)
    except Exception:
        pass
    return result or None


def font_payload(font: Any) -> dict[str, Any]:
    result: dict[str, Any] = {}
    for attr in ("name", "size", "bold", "italic", "underline"):
        try:
            value = getattr(font, attr)
            result[attr] = int(value) if attr == "size" and value is not None else value
        except Exception:
            pass
    try:
        result["color"] = color_payload(font.color)
    except Exception:
        pass
    return result


def text_payload(shape: Any) -> dict[str, Any]:
    if not getattr(shape, "has_text_frame", False):
        return {"has_text_frame": False, "text": "", "paragraphs": []}

    text_frame = shape.text_frame
    paragraphs: list[dict[str, Any]] = []
    for paragraph in text_frame.paragraphs:
        runs = [
            {"text": run.text, "font": font_payload(run.font)}
            for run in paragraph.runs
        ]
        paragraphs.append(
            {
                "text": paragraph.text,
                "alignment": enum_name(paragraph.alignment),
                "level": paragraph.level,
                "runs": runs,
            }
        )

    margins: dict[str, Any] = {}
    for attr in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        try:
            value = getattr(text_frame, attr)
            margins[attr] = int(value) if value is not None else None
        except Exception:
            pass

    return {
        "has_text_frame": True,
        "text": "\n".join(paragraph.text for paragraph in text_frame.paragraphs).strip(),
        "word_wrap": text_frame.word_wrap,
        "auto_size": enum_name(text_frame.auto_size),
        "vertical_anchor": enum_name(text_frame.vertical_anchor),
        "margins": margins,
        "paragraphs": paragraphs,
    }


def placeholder_payload(shape: Any) -> dict[str, Any] | None:
    if not getattr(shape, "is_placeholder", False):
        return None
    placeholder = shape.placeholder_format
    return {
        "idx": placeholder.idx,
        "type": str(placeholder.type),
        "type_name": enum_name(placeholder.type),
    }


def bbox_payload(shape: Any, slide_w: int, slide_h: int) -> dict[str, Any]:
    left = int(shape.left)
    top = int(shape.top)
    width = int(shape.width)
    height = int(shape.height)
    return {
        "emu": {"x": left, "y": top, "width": width, "height": height},
        "inches": {
            "x": emu_to_inches(left),
            "y": emu_to_inches(top),
            "width": emu_to_inches(width),
            "height": emu_to_inches(height),
        },
        "normalized": {
            "x": round(left / slide_w, 6) if slide_w else None,
            "y": round(top / slide_h, 6) if slide_h else None,
            "width": round(width / slide_w, 6) if slide_w else None,
            "height": round(height / slide_h, 6) if slide_h else None,
        },
    }


def slide_relationship_targets(pptx_path: Path, slide_number: int) -> dict[str, str]:
    rels_xml = f"ppt/slides/_rels/slide{slide_number}.xml.rels"
    with ZipFile(pptx_path) as archive:
        if rels_xml not in archive.namelist():
            return {}
        root = ET.fromstring(archive.read(rels_xml))

    rels: dict[str, str] = {}
    for rel in root.findall("rel:Relationship", NS):
        rel_id = rel.attrib.get("Id")
        target = rel.attrib.get("Target")
        if rel_id and target:
            rels[rel_id] = target
    return rels


def picture_relationship_ids(pptx_path: Path, slide_number: int) -> dict[str, str]:
    slide_xml = f"ppt/slides/slide{slide_number}.xml"
    with ZipFile(pptx_path) as archive:
        root = ET.fromstring(archive.read(slide_xml))

    rel_ids: dict[str, str] = {}
    for picture in root.findall(".//p:pic", NS):
        c_nv_pr = picture.find(".//p:cNvPr", NS)
        blip = picture.find(".//a:blip", NS)
        if c_nv_pr is None or blip is None:
            continue
        shape_id = c_nv_pr.attrib.get("id")
        rel_id = blip.attrib.get(f"{{{NS['r']}}}embed") or blip.attrib.get(f"{{{NS['r']}}}link")
        if shape_id and rel_id:
            rel_ids[shape_id] = rel_id
    return rel_ids


def media_payload(pptx_path: Path, relationship_target: str | None) -> dict[str, Any] | None:
    if relationship_target is None:
        return None
    media_path = relationship_target
    if media_path.startswith("../"):
        media_path = f"ppt/{media_path[3:]}"

    with ZipFile(pptx_path) as archive:
        if media_path not in archive.namelist():
            return {
                "relationship_target": relationship_target,
                "resolved_path": media_path,
                "found": False,
            }
        data = archive.read(media_path)

    result: dict[str, Any] = {
        "relationship_target": relationship_target,
        "resolved_path": media_path,
        "found": True,
        "bytes": len(data),
    }
    try:
        with Image.open(io.BytesIO(data)) as image:
            result["image"] = {
                "format": image.format,
                "width_px": image.width,
                "height_px": image.height,
                "mode": image.mode,
            }
    except Exception as exc:
        result["image_error"] = str(exc)
    return result


def picture_payload(
    shape: Any,
    *,
    pptx_path: Path,
    relationship_ids: dict[str, str],
    relationship_targets: dict[str, str],
) -> dict[str, Any] | None:
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return None
    shape_id = str(shape.shape_id)
    rel_id = relationship_ids.get(shape_id)
    target = relationship_targets.get(rel_id or "")
    crop: dict[str, Any] = {}
    for attr in ("crop_left", "crop_right", "crop_top", "crop_bottom"):
        try:
            crop[attr] = getattr(shape, attr)
        except Exception:
            pass
    return {
        "relationship_id": rel_id,
        "relationship_target": target,
        "crop": crop,
        "media": media_payload(pptx_path, target),
    }


def shape_payload(shape: Any) -> dict[str, Any]:
    return {
        "id": str(shape.shape_id),
        "name": shape.name,
        "shape_type": str(shape.shape_type),
        "text": text_payload(shape)["text"],
        "x": emu_to_inches(int(shape.left)),
        "y": emu_to_inches(int(shape.top)),
        "width": emu_to_inches(int(shape.width)),
        "height": emu_to_inches(int(shape.height)),
        "area": int(shape.width) * int(shape.height),
    }


def select_map_shape(
    shapes: list[Any],
    *,
    map_shape_id: str | None,
    map_shape_name: str | None,
) -> Any:
    if map_shape_id is not None:
        for shape in shapes:
            if str(shape.shape_id) == map_shape_id:
                return shape
        raise ValueError(f"Không tìm thấy shape id `{map_shape_id}` trong slide.")

    if map_shape_name is not None:
        for shape in shapes:
            if shape.name == map_shape_name:
                return shape
        raise ValueError(f"Không tìm thấy shape tên `{map_shape_name}` trong slide.")

    pictures = [shape for shape in shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    candidates = pictures or shapes
    return max(candidates, key=lambda shape: int(shape.width) * int(shape.height))


def relative_template_path(pptx_path: Path, output_path: Path) -> str:
    return os.path.relpath(pptx_path, output_path.parent).replace("\\", "/")


def detailed_shape_payload(
    shape: Any,
    *,
    slide_w: int,
    slide_h: int,
    pptx_path: Path,
    relationship_ids: dict[str, str],
    relationship_targets: dict[str, str],
) -> dict[str, Any]:
    payload = {
        "id": str(shape.shape_id),
        "name": shape.name,
        "shape_type": str(shape.shape_type),
        "shape_type_name": enum_name(shape.shape_type),
        "rotation": shape.rotation,
        "bbox": bbox_payload(shape, slide_w, slide_h),
        "placeholder": placeholder_payload(shape),
        "text": text_payload(shape),
        "fill": fill_payload(shape),
        "line": line_payload(shape),
        "picture": picture_payload(
            shape,
            pptx_path=pptx_path,
            relationship_ids=relationship_ids,
            relationship_targets=relationship_targets,
        ),
    }
    return payload


def slide_inspection_payload(
    presentation: Any,
    pptx_path: Path,
    slide_index: int,
) -> dict[str, Any]:
    slide_number = slide_index + 1
    slide = presentation.slides[slide_index]
    slide_w = int(presentation.slide_width)
    slide_h = int(presentation.slide_height)
    relationship_ids = picture_relationship_ids(pptx_path, slide_number)
    relationship_targets = slide_relationship_targets(pptx_path, slide_number)
    return {
        "slide_index": slide_index,
        "slide_number": slide_number,
        "shape_count": len(slide.shapes),
        "relationship_targets": relationship_targets,
        "shapes": [
            detailed_shape_payload(
                shape,
                slide_w=slide_w,
                slide_h=slide_h,
                pptx_path=pptx_path,
                relationship_ids=relationship_ids,
                relationship_targets=relationship_targets,
            )
            for shape in slide.shapes
        ],
    }


def inspection_metadata(
    presentation: Any,
    pptx_path: Path,
    output_path: Path,
    slide_index: int,
    map_shape: Any,
) -> dict[str, Any]:
    slide_w = int(presentation.slide_width)
    slide_h = int(presentation.slide_height)
    slides = [
        slide_inspection_payload(presentation, pptx_path, index)
        for index in range(len(presentation.slides))
    ]
    return {
        "generator": {
            "name": Path(__file__).name,
            "schema": "pptx-template-inspection-v1",
        },
        "source": {
            "pptx": relative_template_path(pptx_path, output_path),
            "absolute_pptx": str(pptx_path),
        },
        "presentation": {
            "slide_count": len(presentation.slides),
            "slide_width_emu": slide_w,
            "slide_height_emu": slide_h,
            "slide_width_inches": emu_to_inches(slide_w),
            "slide_height_inches": emu_to_inches(slide_h),
        },
        "selected_slide": slides[slide_index],
        "slides": slides,
        "map_shape": {
            "id": str(map_shape.shape_id),
            "name": map_shape.name,
            "selection_rule": "explicit option or largest picture on selected slide",
        },
        "limitations": [
            "This JSON is metadata for inspecting and addressing the PPTX template.",
            "It is not a full reversible serialization of the PPTX package.",
            "The original PPTX file is still required for export and editing.",
        ],
    }


def build_metadata(
    pptx_path: Path,
    output_path: Path,
    *,
    slide_index: int,
    map_shape_id: str | None,
    map_shape_name: str | None,
) -> TemplateMetadata:
    presentation = Presentation(str(pptx_path))
    if slide_index < 0 or slide_index >= len(presentation.slides):
        raise ValueError(
            f"slide-index {slide_index} ngoài phạm vi; file có {len(presentation.slides)} slide."
        )

    slide = presentation.slides[slide_index]
    shapes = list(slide.shapes)
    map_shape = select_map_shape(
        shapes,
        map_shape_id=map_shape_id,
        map_shape_name=map_shape_name,
    )
    map_info = shape_payload(map_shape)

    placeholders = [
        TemplatePlaceholder(
            field="map_image",
            element_id=int(map_shape.shape_id),
            kind="map_image",
            diagnostic_name=map_shape.name,
            required=True,
        )
    ]
    for shape in shapes:
        if shape == map_shape or not getattr(shape, "has_text_frame", False):
            continue
        text = shape_payload(shape)["text"]
        if not text:
            continue
        placeholders.append(
            TemplatePlaceholder(
                field=f"text_{shape.shape_id}",
                element_id=int(shape.shape_id),
                kind="text",
                diagnostic_name=shape.name,
                required=False,
            )
        )

    return TemplateMetadata(
        template_pptx=relative_template_path(pptx_path, output_path),
        slide_index=slide_index,
        map_frame=MapFrame(
            x=map_info["x"],
            y=map_info["y"],
            width=map_info["width"],
            height=map_info["height"],
        ),
        placeholders=placeholders,
        metadata=inspection_metadata(
            presentation,
            pptx_path,
            output_path,
            slide_index,
            map_shape,
        ),
    )


def parse_args(argv: list[str]) -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Generate TemplateMetadata JSON from a PPTX file."
    )
    parser.add_argument("pptx", type=Path, help="PowerPoint template file.")
    parser.add_argument(
        "--output",
        type=Path,
        default=Path("examples/templates/target_001.template.json"),
        help="Output metadata JSON path.",
    )
    parser.add_argument("--slide-index", type=int, default=0, help="Zero-based slide index.")
    parser.add_argument("--map-shape-id", help="Explicit map image shape id.")
    parser.add_argument("--map-shape-name", help="Explicit map image shape name.")
    return parser.parse_args(argv)


def main(argv: list[str]) -> int:
    args = parse_args(argv)
    pptx_path = args.pptx.resolve()
    output_path = args.output.resolve()
    if not pptx_path.is_file():
        print(f"ERROR: Không tìm thấy PPTX: {pptx_path}", file=sys.stderr)
        return 2

    metadata = build_metadata(
        pptx_path,
        output_path,
        slide_index=args.slide_index,
        map_shape_id=args.map_shape_id,
        map_shape_name=args.map_shape_name,
    )
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(
        json.dumps(metadata.model_dump(mode="json"), ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )
    print(f"Wrote {output_path}")
    print(
        "Map frame:",
        metadata.placeholders[0].diagnostic_name,
        f"id={metadata.placeholders[0].element_id}",
        metadata.map_frame.model_dump(),
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
