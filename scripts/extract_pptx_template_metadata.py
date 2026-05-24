#!/usr/bin/env python3
"""Extract PowerPoint slide template metadata for 3.ThucTheNgay.

The output JSON is intended to be used as target-specific template metadata.
It records slide size, shape names/ids, bounding boxes, text placeholders,
picture/image information, and a best-effort map-frame candidate.
"""

from __future__ import annotations

import argparse
import io
import json
import re
import sys
from dataclasses import dataclass
from pathlib import Path
from typing import Any
from zipfile import ZipFile

from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


EMU_PER_INCH = 914400
NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
}


@dataclass(frozen=True)
class XmlShapeInfo:
    shape_id: int
    name: str
    title: str | None
    descr: str | None
    rel_id: str | None
    xml_tag: str


def emu_to_inches(value: int) -> float:
    return round(value / EMU_PER_INCH, 4)


def bbox(left: int, top: int, width: int, height: int, slide_w: int, slide_h: int) -> dict[str, Any]:
    return {
        "emu": {"x": left, "y": top, "width": width, "height": height},
        "inches": {
            "x": emu_to_inches(left),
            "y": emu_to_inches(top),
            "width": emu_to_inches(width),
            "height": emu_to_inches(height),
        },
        "normalized": {
            "x": round(left / slide_w, 6) if slide_w else None,
            "y": round(top / slide_h, 6) if slide_h else None,
            "width": round(width / slide_w, 6) if slide_w else None,
            "height": round(height / slide_h, 6) if slide_h else None,
        },
    }


def local_name(tag: str) -> str:
    return etree.QName(tag).localname


def read_xml_shape_info(pptx_path: Path, slide_number: int) -> tuple[dict[int, XmlShapeInfo], dict[str, str]]:
    slide_xml = f"ppt/slides/slide{slide_number}.xml"
    rels_xml = f"ppt/slides/_rels/slide{slide_number}.xml.rels"
    with ZipFile(pptx_path) as zf:
        slide_root = etree.fromstring(zf.read(slide_xml))
        rel_root = etree.fromstring(zf.read(rels_xml)) if rels_xml in zf.namelist() else None

    rel_targets: dict[str, str] = {}
    if rel_root is not None:
        for rel in rel_root.xpath("./rel:Relationship", namespaces=NS):
            rel_id = rel.get("Id")
            target = rel.get("Target")
            if rel_id and target:
                rel_targets[rel_id] = target

    info: dict[int, XmlShapeInfo] = {}
    for node in slide_root.xpath(".//*[p:cNvPr]", namespaces=NS):
        c_nv_pr_nodes = node.xpath("./*/p:cNvPr", namespaces=NS)
        if not c_nv_pr_nodes:
            continue
        c_nv_pr = c_nv_pr_nodes[0]
        raw_id = c_nv_pr.get("id")
        if raw_id is None:
            continue
        rel_id = None
        blip_nodes = node.xpath(".//a:blip", namespaces=NS)
        if blip_nodes:
            rel_id = blip_nodes[0].get(f"{{{NS['r']}}}embed") or blip_nodes[0].get(f"{{{NS['r']}}}link")
        shape_id = int(raw_id)
        info[shape_id] = XmlShapeInfo(
            shape_id=shape_id,
            name=c_nv_pr.get("name") or "",
            title=c_nv_pr.get("title"),
            descr=c_nv_pr.get("descr"),
            rel_id=rel_id,
            xml_tag=local_name(node.tag),
        )
    return info, rel_targets


def media_info(pptx_path: Path, target: str | None) -> dict[str, Any] | None:
    if not target:
        return None
    normalized = target
    if normalized.startswith("../"):
        normalized = "ppt/" + normalized[3:]
    with ZipFile(pptx_path) as zf:
        if normalized not in zf.namelist():
            return {"target": target, "resolved_path": normalized, "found": False}
        data = zf.read(normalized)
    result: dict[str, Any] = {
        "target": target,
        "resolved_path": normalized,
        "found": True,
        "bytes": len(data),
    }
    try:
        with Image.open(io.BytesIO(data)) as image:
            result["image"] = {
                "format": image.format,
                "width_px": image.width,
                "height_px": image.height,
                "mode": image.mode,
            }
    except Exception as exc:
        result["image_error"] = str(exc)
    return result


def shape_type_name(shape: Any) -> str:
    try:
        return MSO_SHAPE_TYPE(shape.shape_type).name
    except Exception:
        return str(shape.shape_type)


def placeholder_info(shape: Any) -> dict[str, Any] | None:
    if not getattr(shape, "is_placeholder", False):
        return None
    fmt = shape.placeholder_format
    return {
        "idx": fmt.idx,
        "type": str(fmt.type),
        "type_name": getattr(fmt.type, "name", str(fmt.type)),
    }


def extract_text(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs).strip()


def fill_info(shape: Any) -> dict[str, Any] | None:
    fill = getattr(shape, "fill", None)
    if fill is None:
        return None
    result: dict[str, Any] = {"type": str(fill.type)}
    try:
        if fill.fore_color and fill.fore_color.rgb:
            result["rgb"] = str(fill.fore_color.rgb)
    except Exception:
        pass
    return result


def line_info(shape: Any) -> dict[str, Any] | None:
    line = getattr(shape, "line", None)
    if line is None:
        return None
    result: dict[str, Any] = {}
    try:
        result["width_emu"] = int(line.width) if line.width is not None else None
    except Exception:
        pass
    try:
        if line.color and line.color.rgb:
            result["rgb"] = str(line.color.rgb)
    except Exception:
        pass
    return result or None


def is_text_placeholder_candidate(text: str) -> bool:
    if not text:
        return False
    patterns = [
        r"\btarget\b",
        r"\bhello world\b",
        r"\d{1,2}[.:]\d{2}/\d{1,2}\.\d{1,2}\.\d{2,4}",
        r"ngày\s+\d{1,2}\.\d{1,2}\.\d{2,4}",
        r"xyz",
    ]
    lower = text.lower()
    return any(re.search(pattern, lower, flags=re.IGNORECASE) for pattern in patterns)


def classify_shape(shape: dict[str, Any], slide_area: int) -> list[str]:
    tags: list[str] = []
    text = shape.get("text") or ""
    if shape["type"] == "PICTURE":
        tags.append("image")
        if shape["area_emu"] >= slide_area * 0.25:
            tags.append("map_frame_candidate")
    if text:
        tags.append("text")
        if is_text_placeholder_candidate(text):
            tags.append("text_placeholder_candidate")
    if shape.get("placeholder"):
        tags.append("powerpoint_placeholder")
    return tags


def extract_metadata(pptx_path: Path, slide_index: int) -> dict[str, Any]:
    prs = Presentation(str(pptx_path))
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"slide-index {slide_index} out of range; presentation has {len(prs.slides)} slide(s)")

    slide_number = slide_index + 1
    slide = prs.slides[slide_index]
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)
    slide_area = slide_w * slide_h
    xml_info, rel_targets = read_xml_shape_info(pptx_path, slide_number)

    shapes: list[dict[str, Any]] = []
    for shape in slide.shapes:
        shape_id = int(shape.shape_id)
        xinfo = xml_info.get(shape_id)
        rel_target = rel_targets.get(xinfo.rel_id) if xinfo and xinfo.rel_id else None
        item: dict[str, Any] = {
            "id": shape_id,
            "name": shape.name,
            "xml_name": xinfo.name if xinfo else None,
            "title": xinfo.title if xinfo else None,
            "descr": xinfo.descr if xinfo else None,
            "xml_tag": xinfo.xml_tag if xinfo else None,
            "type": shape_type_name(shape),
            "placeholder": placeholder_info(shape),
            "text": extract_text(shape),
            "bbox": bbox(int(shape.left), int(shape.top), int(shape.width), int(shape.height), slide_w, slide_h),
            "area_emu": int(shape.width) * int(shape.height),
            "fill": fill_info(shape),
            "line": line_info(shape),
            "relationship_id": xinfo.rel_id if xinfo else None,
            "relationship_target": rel_target,
            "media": media_info(pptx_path, rel_target),
        }
        item["tags"] = classify_shape(item, slide_area)
        shapes.append(item)

    map_candidates = sorted(
        [shape for shape in shapes if "map_frame_candidate" in shape["tags"]],
        key=lambda item: item["area_emu"],
        reverse=True,
    )
    text_candidates = [shape for shape in shapes if "text_placeholder_candidate" in shape["tags"]]
    placeholders = {
        "map_image": {
            "shape_name": map_candidates[0]["name"],
            "fallback_id": map_candidates[0]["id"],
            "bbox": map_candidates[0]["bbox"],
            "reason": "largest picture occupying at least 25% of slide area",
        }
        if map_candidates
        else None,
        "text_candidates": [
            {
                "shape_name": shape["name"],
                "fallback_id": shape["id"],
                "current_text": shape["text"],
                "bbox": shape["bbox"],
            }
            for shape in text_candidates
        ],
    }

    return {
        "schema_version": "template-metadata-v0",
        "template_pptx": str(pptx_path),
        "slide_index": slide_index,
        "slide_number": slide_number,
        "presentation": {
            "slide_count": len(prs.slides),
            "slide_width_emu": slide_w,
            "slide_height_emu": slide_h,
            "slide_width_inches": emu_to_inches(slide_w),
            "slide_height_inches": emu_to_inches(slide_h),
        },
        "placeholders": placeholders,
        "shapes": shapes,
    }


def parse_args(argv: list[str]) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Extract PPTX template metadata to JSON.")
    parser.add_argument("pptx", type=Path, help="Path to PowerPoint template file.")
    parser.add_argument("--slide-index", type=int, default=0, help="Zero-based slide index to inspect. Default: 0.")
    parser.add_argument("--output", type=Path, help="Output JSON path. Default: <pptx-stem>.template.json")
    parser.add_argument("--pretty", action="store_true", help="Pretty-print JSON with indentation.")
    return parser.parse_args(argv)


def main(argv: list[str]) -> int:
    args = parse_args(argv)
    pptx_path = args.pptx.resolve()
    if not pptx_path.exists():
        print(f"ERROR: PPTX file not found: {pptx_path}", file=sys.stderr)
        return 2
    output = args.output or pptx_path.with_suffix(".template.json")
    metadata = extract_metadata(pptx_path, args.slide_index)
    output.write_text(
        json.dumps(metadata, ensure_ascii=False, indent=2 if args.pretty else None) + "\n",
        encoding="utf-8",
    )
    print(f"Wrote {output}")
    map_image = metadata["placeholders"]["map_image"]
    if map_image:
        print(f"Map image candidate: {map_image['shape_name']} (id={map_image['fallback_id']})")
    print(f"Shapes: {len(metadata['shapes'])}")
    print(f"Text placeholder candidates: {len(metadata['placeholders']['text_candidates'])}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
