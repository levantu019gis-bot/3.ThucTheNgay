"""Low-level PowerPoint slide copy and element-id replacement helpers."""

from __future__ import annotations

from copy import deepcopy
from io import BytesIO
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.slide import Slide


def copy_only_slide(source: Presentation, destination: Presentation) -> Slide:
    """Copy the only slide from ``source`` into ``destination``."""
    if len(source.slides) != 1:
        msg = f"Expected one template slide, got {len(source.slides)}"
        raise ValueError(msg)

    slide = destination.slides.add_slide(destination.slide_layouts[6])
    source_slide = source.slides[0]
    for shape in source_slide.shapes:
        slide.shapes._spTree.insert_element_before(deepcopy(shape.element), "p:extLst")

    rel_id_map: dict[str, str] = {}
    for rel in source_slide.part.rels.values():
        if rel.is_external:
            continue
        if "notesSlide" in rel.reltype:
            continue
        if rel.rId in slide.part.rels:
            continue
        if rel.reltype.endswith("/image"):
            _image_part, rel_id_map[rel.rId] = slide.part.get_or_add_image_part(
                BytesIO(rel._target.blob)
            )
        else:
            rel_id_map[rel.rId] = slide.part.rels._add_relationship(rel.reltype, rel._target)
    _remap_relationship_ids(slide, rel_id_map)
    return slide


def find_shape_by_element_id(slide: Slide, element_id: int) -> Any | None:
    """Return the shape matching a PowerPoint element id, recursing into groups."""
    return _find_shape_by_element_id(slide.shapes, element_id)


def replace_shape_with_picture(slide: Slide, element_id: int, image_path: Path) -> bool:
    """Replace a shape rectangle with a picture at the same coordinates."""
    shape = find_shape_by_element_id(slide, element_id)
    if shape is None:
        return False
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    shape.element.getparent().remove(shape.element)
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    return True


def replace_text(slide: Slide, element_id: int, value: str) -> bool:
    """Replace the text of a shape addressed by element id."""
    shape = find_shape_by_element_id(slide, element_id)
    if shape is None or not getattr(shape, "has_text_frame", False):
        return False
    shape.text = value
    return True


def _find_shape_by_element_id(shapes: Any, element_id: int) -> Any | None:
    for shape in shapes:
        if int(shape.shape_id) == element_id:
            return shape
        child_shapes = getattr(shape, "shapes", None)
        if child_shapes is not None:
            found = _find_shape_by_element_id(child_shapes, element_id)
            if found is not None:
                return found
    return None


def _remap_relationship_ids(slide: Slide, rel_id_map: dict[str, str]) -> None:
    if not rel_id_map:
        return
    for element in slide.element.iter():
        for attr_name, attr_value in tuple(element.attrib.items()):
            replacement = rel_id_map.get(attr_value)
            if replacement is not None:
                element.set(attr_name, replacement)
