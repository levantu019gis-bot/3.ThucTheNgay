"""Load and validate target-specific one-slide PowerPoint templates."""

from __future__ import annotations

from collections.abc import Iterable
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation
from pydantic import ValidationError

from thucthengay.models import (
    Issue,
    IssueScope,
    IssueSeverity,
    MapFrame,
    PlaceholderType,
    TargetConfig,
    TemplateMetadata,
    TemplatePlaceholder,
)

EMU_PER_POINT = 12_700.0


class TemplateLoadError(Exception):
    """Expected template loading failure with issue details."""

    def __init__(self, issue_id: str, message: str, remediation: str) -> None:
        self.issue_id = issue_id
        self.message = message
        self.remediation = remediation
        super().__init__(message)


@dataclass(frozen=True)
class LoadedTemplate:
    """Resolved template information needed by render/export preflight."""

    target_id: str
    template_pptx_file: Path
    metadata: TemplateMetadata
    element_names: dict[int, str]
    compatibility_signature: str


def load_target_template(target: TargetConfig, template_pptx_file: Path) -> LoadedTemplate:
    """Load one target PPTX template and verify configured element-id mappings."""
    if not template_pptx_file.is_file():
        raise TemplateLoadError(
            "target.template_pptx_missing",
            f"Khong tim thay PPTX template cua target `{target.id}`: {template_pptx_file}",
            "Kiem tra lai `export.template_pptx_file` trong config; duong dan tinh tu config.json.",
        )

    try:
        presentation = Presentation(str(template_pptx_file))
    except Exception as exc:
        raise TemplateLoadError(
            "target.template_pptx_invalid",
            f"PPTX template cua target `{target.id}` khong doc duoc: {template_pptx_file}",
            f"Mo lai file PPTX va luu dung dinh dang PowerPoint. Chi tiet ky thuat: {exc}",
        ) from exc

    slide_count = len(presentation.slides)
    if slide_count != 1:
        raise TemplateLoadError(
            "target.template_pptx_slide_count_invalid",
            f"PPTX template cua target `{target.id}` phai co dung 1 slide, hien co {slide_count}.",
            "Tach moi target thanh mot file PPTX template rieng chi gom 1 slide mau.",
        )

    slide = presentation.slides[0]
    shapes_by_id = _shapes_by_element_id(slide.shapes)
    element_names = {
        element_id: getattr(shape, "name", "")
        for element_id, shape in shapes_by_id.items()
    }
    placeholders = _placeholders_with_diagnostics(target.export.placeholders, element_names)
    _reject_duplicate_element_ids(placeholders, target.id)
    map_placeholder = _required_map_placeholder(placeholders, target.id)

    for placeholder in placeholders:
        if placeholder.required and placeholder.element_id not in shapes_by_id:
            raise TemplateLoadError(
                "target.template_element_missing",
                (
                    f"PPTX template cua target `{target.id}` thieu element id "
                    f"`{placeholder.element_id}` cho field `{placeholder.field}`."
                ),
                (
                    "Cap nhat `export.placeholders[].element_id` theo dung shape id "
                    "trong PPTX template."
                ),
            )

    map_shape = shapes_by_id[map_placeholder.element_id]
    try:
        map_frame = MapFrame(
            x=_emu_to_points(map_shape.left),
            y=_emu_to_points(map_shape.top),
            width=_emu_to_points(map_shape.width),
            height=_emu_to_points(map_shape.height),
        )
    except ValidationError as exc:
        raise TemplateLoadError(
            "target.template_map_frame_invalid",
            (
                f"Khung ban do trong PPTX template cua target `{target.id}` khong hop le "
                f"o element id `{map_placeholder.element_id}`."
            ),
            (
                "Sua vi tri/kich thuoc shape map image trong PPTX de x/y khong am "
                "va width/height lon hon 0."
            ),
        ) from exc

    return LoadedTemplate(
        target_id=target.id,
        template_pptx_file=template_pptx_file,
        metadata=TemplateMetadata(
            template_pptx=str(template_pptx_file),
            slide_index=0,
            map_frame=map_frame,
            placeholders=placeholders,
            metadata={
                "source": "template_pptx_file",
                "element_names": {str(key): value for key, value in element_names.items()},
            },
        ),
        element_names=element_names,
        compatibility_signature=_compatibility_signature(presentation),
    )


def template_compatibility_issues(templates: Iterable[LoadedTemplate]) -> tuple[Issue, ...]:
    """Surface detectable or unknown compatibility risks across PPTX templates."""
    loaded = tuple(templates)
    files = {template.template_pptx_file for template in loaded}
    if len(files) <= 1:
        return ()

    signatures = {template.compatibility_signature for template in loaded}
    detail = (
        "khac nhau"
        if len(signatures) > 1
        else "chua xac minh chac chan du thong tin master/layout doc duoc giong nhau"
    )
    target_ids = ", ".join(template.target_id for template in loaded)
    return (
        Issue(
            issue_id="target.template_compatibility_unknown",
            severity=IssueSeverity.WARNING,
            scope=IssueScope.TEMPLATE,
            message=f"Nhieu PPTX template co base/theme/master {detail}.",
            remediation=(
                "Dam bao cac template cua target duoc tao tu cung mot PowerPoint base/theme/master "
                f"truoc khi export tong hop. Targets: {target_ids}."
            ),
        ),
    )


def _shapes_by_element_id(shapes: Any) -> dict[int, Any]:
    result: dict[int, Any] = {}
    for shape in shapes:
        result[int(shape.shape_id)] = shape
        child_shapes = getattr(shape, "shapes", None)
        if child_shapes is not None:
            result.update(_shapes_by_element_id(child_shapes))
    return result


def _placeholders_with_diagnostics(
    placeholders: Iterable[TemplatePlaceholder],
    element_names: dict[int, str],
) -> list[TemplatePlaceholder]:
    return [
        placeholder.model_copy(
            update={
                "diagnostic_name": placeholder.diagnostic_name
                or element_names.get(placeholder.element_id)
            }
        )
        for placeholder in placeholders
    ]


def _reject_duplicate_element_ids(
    placeholders: Iterable[TemplatePlaceholder],
    target_id: str,
) -> None:
    seen: dict[int, str] = {}
    for placeholder in placeholders:
        prior_field = seen.get(placeholder.element_id)
        if prior_field is not None:
            raise TemplateLoadError(
                "target.template_element_duplicate",
                (
                    f"Target `{target_id}` cau hinh trung element id "
                    f"`{placeholder.element_id}` cho field `{prior_field}` "
                    f"va `{placeholder.field}`."
                ),
                "Moi placeholder export phai tro toi mot element id rieng trong PPTX template.",
            )
        seen[placeholder.element_id] = placeholder.field


def _required_map_placeholder(
    placeholders: Iterable[TemplatePlaceholder],
    target_id: str,
) -> TemplatePlaceholder:
    matches = [
        placeholder
        for placeholder in placeholders
        if placeholder.kind == PlaceholderType.MAP_IMAGE and placeholder.required
    ]
    if len(matches) > 1:
        fields = ", ".join(placeholder.field for placeholder in matches)
        raise TemplateLoadError(
            "target.template_map_element_ambiguous",
            f"Target `{target_id}` co nhieu map image placeholder bat buoc: {fields}.",
            (
                "Chi giu mot `kind: map_image` bat buoc lam khung ban do chinh "
                "trong `export.placeholders`."
            ),
        )
    if matches:
        return matches[0]
    raise TemplateLoadError(
        "target.template_map_element_missing",
        f"Target `{target_id}` chua khai bao map image placeholder bat buoc.",
        (
            "Bo sung `export.placeholders` voi `kind: map_image`, `field`, "
            "va `element_id` cua khung ban do."
        ),
    )


def _emu_to_points(value: int) -> float:
    return float(value) / EMU_PER_POINT


def _compatibility_signature(presentation: Presentation) -> str:
    masters = getattr(presentation.slide_masters, "_sldMasterIdLst", None)
    master_count = len(presentation.slide_masters)
    layouts_count = len(presentation.slide_layouts)
    master_xml = masters.xml if masters is not None else ""
    return f"masters={master_count};layouts={layouts_count};master_ids={master_xml}"
